import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# A constant to convert inches (used in the JSON) to EMUs (the internal unit for pptx)
EMUS_PER_INCH = 914400

# The JSON data provided by the user
diagram_json = {
    "canvas": {
        "width": 13.333,
        "height": 7.5
    },
    "containers": {
        "web-ui": {
            "left": 1.4249875,
            "top": 1.4125,
            "width": 1.416625,
            "height": 4.675
        },
        "api-gateway": {
            "left": 3.2666,
            "top": 1.4125,
            "width": 1.416625,
            "height": 1.44375
        },
        "domain-logic": {
            "left": 3.2666,
            "top": 2.99375,
            "width": 1.416625,
            "height": 1.44375
        },
        "order-queue": {
            "left": 3.2666,
            "top": 4.575,
            "width": 1.416625,
            "height": 1.44375
        },
        "application-services": {
            "left": 5.1082125,
            "top": 1.4125,
            "width": 1.416625,
            "height": 4.675
        },
        "persistence-layer": {
            "left": 6.949825,
            "top": 1.4125,
            "width": 1.416625,
            "height": 4.675
        },
        "integration-layer": {
            "left": 8.7914375,
            "top": 1.4125,
            "width": 1.416625,
            "height": 4.675
        },
        "security": {
            "left": 10.63305,
            "top": 1.4125,
            "width": 1.416625,
            "height": 1.44375
        },
        "platform": {
            "left": 10.63305,
            "top": 2.99375,
            "width": 1.416625,
            "height": 1.44375
        },
        "deployment": {
            "left": 10.63305,
            "top": 4.575,
            "width": 1.416625,
            "height": 1.44375
        }
    },
    "arrows": [
        {
            "id": "user-request",
            "from": "web-ui",
            "to": "api-gateway",
            "path": [
                {
                    "x": 2.91244375,
                    "y": 2.134375
                },
                {
                    "x": 3.19576875,
                    "y": 2.134375
                }
            ],
            "label": {
                "text": "request",
                "x": 3.05410625,
                "y": 2.1
            }
        },
        {
            "id": "authentication",
            "from": "api-gateway",
            "to": "security",
            "path": [
                {
                    "x": 4.683225,
                    "y": 1.89375
                },
                {
                    "x": 4.89571875,
                    "y": 1.89375
                },
                {
                    "x": 4.89571875,
                    "y": 1.34375
                },
                {
                    "x": 10.42055625,
                    "y": 1.34375
                },
                {
                    "x": 10.42055625,
                    "y": 2.134375
                },
                {
                    "x": 10.56221875,
                    "y": 2.134375
                }
            ],
            "label": {
                "text": "auth",
                "x": 7.6581375,
                "y": 1.309375
            }
        },
        {
            "id": "business-logic",
            "from": "api-gateway",
            "to": "domain-logic",
            "path": [
                {
                    "x": 3.9749125,
                    "y": 2.87
                },
                {
                    "x": 3.9749125,
                    "y": 2.98
                }
            ],
            "label": {
                "text": "process",
                "x": 4.081159375,
                "y": 2.925
            }
        },
        {
            "id": "service-execution",
            "from": "domain-logic",
            "to": "application-services",
            "path": [
                {
                    "x": 4.75405625,
                    "y": 3.715625
                },
                {
                    "x": 5.03738125,
                    "y": 3.715625
                }
            ],
            "label": {
                "text": "execute",
                "x": 4.89571875,
                "y": 3.68125
            }
        },
        {
            "id": "data-operations",
            "from": "application-services",
            "to": "persistence-layer",
            "path": [
                {
                    "x": 6.59566875,
                    "y": 3.75
                },
                {
                    "x": 6.87899375,
                    "y": 3.75
                }
            ],
            "label": {
                "text": "data",
                "x": 6.73733125,
                "y": 3.715625
            }
        },
        {
            "id": "async-processing",
            "from": "domain-logic",
            "to": "order-queue",
            "path": [
                {
                    "x": 3.9749125,
                    "y": 4.45125
                },
                {
                    "x": 3.9749125,
                    "y": 4.561249999999999
                }
            ],
            "label": {
                "text": "async",
                "x": 4.081159375,
                "y": 4.50625
            }
        },
        {
            "id": "queue-processing",
            "from": "order-queue",
            "to": "application-services",
            "path": [
                {
                    "x": 4.75405625,
                    "y": 5.296875
                },
                {
                    "x": 5.03738125,
                    "y": 5.296875
                }
            ],
            "label": {
                "text": "process",
                "x": 4.89571875,
                "y": 5.2625
            }
        },
        {
            "id": "external-integration",
            "from": "application-services",
            "to": "integration-layer",
            "path": [
                {
                    "x": 5.816525,
                    "y": 1.48125
                },
                {
                    "x": 5.816525,
                    "y": 1.275
                },
                {
                    "x": 9.49975,
                    "y": 1.275
                },
                {
                    "x": 9.49975,
                    "y": 1.48125
                }
            ],
            "label": {
                "text": "external",
                "x": 7.6581375,
                "y": 1.240625
            }
        },
        {
            "id": "platform-deployment",
            "from": [
                "platform",
                "deployment"
            ],
            "to": "application-services",
            "path": [
                {
                    "x": 10.56221875,
                    "y": 3.715625,
                    "action": "line"
                },
                {
                    "x": 10.349725,
                    "y": 3.715625,
                    "action": "line"
                },
                {
                    "x": 10.56221875,
                    "y": 5.296875,
                    "action": "line"
                },
                {
                    "x": 10.349725,
                    "y": 5.296875,
                    "action": "line"
                },
                {
                    "x": 10.349725,
                    "y": 3.715625,
                    "action": "merge_vertical"
                },
                {
                    "x": 10.349725,
                    "y": 5.296875,
                    "action": "merge_vertical"
                },
                {
                    "x": 10.349725,
                    "y": 4.50625,
                    "action": "branch"
                },
                {
                    "x": 10.349725,
                    "y": 6.29375,
                    "action": "line"
                },
                {
                    "x": 5.816525,
                    "y": 6.29375,
                    "action": "line"
                },
                {
                    "x": 5.816525,
                    "y": 6.0875,
                    "action": "arrow"
                }
            ],
            "label": {
                "text": "monitor & deploy",
                "x": 8.083124999999999,
                "y": 6.259375
            }
        }
    ],
    "arrowhead": {
        "width": 5,
        "height": 4,
        "refX": 4,
        "refY": 2
    }
}


def _add_arrowhead(slide, end_x, end_y, start_x, start_y):
    """
    Manually adds a triangular arrowhead at the end of a line segment.

    Args:
        slide (pptx.slide.Slide): The slide to add the shape to.
        end_x (int): The x-coordinate (in EMUs) of the line's endpoint.
        end_y (int): The y-coordinate (in EMUs) of the line's endpoint.
        start_x (int): The x-coordinate (in EMUs) of the line's startpoint.
        start_y (int): The y-coordinate (in EMUs) of the line's startpoint.
    """
    # Define arrowhead size
    width_in = 0.15 # inches
    height_in = 0.15 # inches

    width_emu = int(width_in * EMUS_PER_INCH)
    height_emu = int(height_in * EMUS_PER_INCH)

    # Calculate the angle of the line
    dx = end_x - start_x
    dy = end_y - start_y
    angle = math.degrees(math.atan2(dy, dx))

    # Calculate the position for the arrowhead
    # We want it to be centered on the end point of the line
    arrow_x = end_x - width_emu / 2
    arrow_y = end_y - height_emu / 2

    arrowhead = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        arrow_x, arrow_y, width_emu, height_emu
    )

    # Rotate the arrowhead to align with the line
    arrowhead.rotation = angle + 90
    
    # Set the arrowhead's color to match the line
    fill = arrowhead.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(105, 105, 105) # Dim gray
    
    # Set the arrowhead's line to be the same color as the fill
    line = arrowhead.line
    line.color.rgb = RGBColor(105, 105, 105)
    line.width = Pt(0.5)

def create_diagram(data):
    """
    Generates a PowerPoint slide with a diagram based on the provided JSON data.

    Args:
        data (dict): A dictionary containing the diagram's structure, coordinates,
                     and connections.
    """
    try:
        # Create a new presentation
        prs = Presentation()
        
        # Set the slide dimensions from the JSON
        slide_width_emu = int(data['canvas']['width'] * EMUS_PER_INCH)
        slide_height_emu = int(data['canvas']['height'] * EMUS_PER_INCH)
        prs.slide_width = slide_width_emu
        prs.slide_height = slide_height_emu

        # Use a blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # A dictionary to store the created shapes for easy lookup
        shapes_map = {}

        # 1. Create the container shapes
        for container_name, props in data["containers"].items():
            left_emu = int(props['left'] * EMUS_PER_INCH)
            top_emu = int(props['top'] * EMUS_PER_INCH)
            width_emu = int(props['width'] * EMUS_PER_INCH)
            height_emu = int(props['height'] * EMUS_PER_INCH)

            # Create a rounded rectangle shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left_emu, top_emu, width_emu, height_emu
            )
            shapes_map[container_name] = shape
            
            # Format the shape
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(176, 196, 222) # Light steel blue
            line = shape.line
            line.color.rgb = RGBColor(105, 105, 105)  # Dim gray
            line.width = Pt(1.5)

            # Add and format the text
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = container_name.replace('-', ' ').title() # Format text for display
            p.font.size = Pt(14)
            p.font.name = 'Arial'
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Center the text vertically
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 2. Create the arrows (lines) and labels
        for arrow in data["arrows"]:
            path_points = arrow['path']
            
            # Draw line segments
            for i in range(len(path_points) - 1):
                start_x_emu = int(path_points[i]['x'] * EMUS_PER_INCH)
                start_y_emu = int(path_points[i]['y'] * EMUS_PER_INCH)
                end_x_emu = int(path_points[i+1]['x'] * EMUS_PER_INCH)
                end_y_emu = int(path_points[i+1]['y'] * EMUS_PER_INCH)

                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    start_x_emu, start_y_emu, end_x_emu, end_y_emu
                )
                
                # Style the line
                line_color = RGBColor(105, 105, 105) # Dim gray
                line.line.width = Pt(1.5)
                line.line.color.rgb = line_color

            # Add the arrowhead to the last line segment
            last_point = path_points[-1]
            second_to_last_point = path_points[-2]
            _add_arrowhead(
                slide,
                int(last_point['x'] * EMUS_PER_INCH),
                int(last_point['y'] * EMUS_PER_INCH),
                int(second_to_last_point['x'] * EMUS_PER_INCH),
                int(second_to_last_point['y'] * EMUS_PER_INCH)
            )
            
            # Add the arrow label as a text box
            label_props = arrow['label']
            label_x_emu = int(label_props['x'] * EMUS_PER_INCH)
            label_y_emu = int(label_props['y'] * EMUS_PER_INCH)
            
            # Estimate label width based on text length
            label_width_emu = len(label_props['text']) * 10000
            
            textbox = slide.shapes.add_textbox(
                label_x_emu, label_y_emu, label_width_emu, Pt(12)
            )
            
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = label_props['text']
            p.font.size = Pt(12)
            p.font.name = 'Arial'
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
        # Save the presentation
        prs.save('diagramx.pptx')
        print("Diagram saved to diagram.pptx")

    except Exception as e:
        print(f"An error occurred: {e}")

if __name__ == '__main__':
    create_diagram(diagram_json)
